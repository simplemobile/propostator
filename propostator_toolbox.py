import streamlit as st
import pandas as pd
from pptx import Presentation
import os
import sqlite3
import hashlib
import shutil
from datetime import datetime
from editar_taxas import taxas_ed


# ======== BANCO HISTÓRICO ========
def criar_banco():
    conn = sqlite3.connect("historico.db")
    cursor = conn.cursor()
    cursor.execute("""
        CREATE TABLE IF NOT EXISTS historico_propostas (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            usuario_nome TEXT,
            usuario_email TEXT,
            usuario_posicao TEXT,
            datahora TEXT,
            tipo TEXT,
            comercial TEXT,
            email TEXT,
            cliente TEXT,
            genero_cliente TEXT,
            opp TEXT,
            frentes TEXT,
            baseline TEXT,
            data_proposta TEXT,
            objetivo TEXT,
            valor_total TEXT,
            horas_mensais TEXT,
            tera_baseline TEXT,
            nivel TEXT
        )
    """)
    conn.commit()
    conn.close()

def salvar_proposta_log(
    usuario, tipo, comercial, email, cliente, genero_cliente, opp,
    frentes, baseline, data, objetivo, valor_total, horas_mensais, tera_baseline, nivel
):
    conn = sqlite3.connect("historico.db")
    cursor = conn.cursor()
    cursor.execute("""
        INSERT INTO historico_propostas (
            usuario_nome, usuario_email, usuario_posicao, datahora, tipo,
            comercial, email, cliente, genero_cliente, opp, frentes, baseline,
            data_proposta, objetivo, valor_total, horas_mensais, tera_baseline, nivel
        ) VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?)
    """, (
        usuario.nome if usuario else None,
        usuario.email if usuario else None,
        usuario.posicao if usuario else None,
        datetime.now().strftime('%Y-%m-%d %H:%M:%S'),
        tipo, comercial, email, cliente, genero_cliente, opp, frentes, baseline,
        data, objetivo, valor_total, horas_mensais, tera_baseline, nivel
    ))
    conn.commit()
    conn.close()

def mostrar_historico(usuario=None):
    conn = sqlite3.connect("historico.db")
    cursor = conn.cursor()
    if usuario:
        cursor.execute("""
            SELECT datahora, tipo, cliente, opp, baseline, horas_mensais, valor_total
            FROM historico_propostas
            WHERE usuario_email = ?
            ORDER BY datahora DESC
            LIMIT 20
        """, (usuario.email,))
    else:
        cursor.execute("""
            SELECT datahora, tipo, cliente, opp, baseline, horas_mensais, valor_total
            FROM historico_propostas
            ORDER BY datahora DESC
            LIMIT 20
        """)
    rows = cursor.fetchall()
    conn.close()
    if rows:
        st.subheader("Histórico de Propostas")
        st.dataframe(
            [
                {
                    "Data/Hora": r[0],
                    "Tipo": r[1],
                    "Cliente": r[2],
                    "OPP": r[3],
                    "Baseline": r[4],
                    "Horas Mensais": r[5],
                    "Valor Total": r[6]
                } for r in rows
            ]
        )
    else:
        st.info("Nenhuma proposta gerada ainda.")

def tela_editar_taxas():
    st.header("Edição de Taxas (Administrador)")

#========= DEFs PARA GERIR TAXAS ========== INICIO =====#
# ---- BANCO RATE_ALOC ----
def get_alocacoes():
    conn = sqlite3.connect("rate_aloc.db")
    cursor = conn.cursor()
    cursor.execute("SELECT nivel, taxa FROM rate_aloc ORDER BY nivel")
    rows = cursor.fetchall()
    conn.close()
    return rows
def update_aloc(nivel, taxa):
    conn = sqlite3.connect("rate_aloc.db")
    cursor = conn.cursor()
    cursor.execute("UPDATE rate_aloc SET taxa = ? WHERE nivel = ?", (taxa, nivel))
    conn.commit()
    conn.close()
# ---- BANCO TAXAS_SERVICO ----
def get_servicos():
    conn = sqlite3.connect("taxas_servico.db")
    cursor = conn.cursor()
    cursor.execute("SELECT servico, taxa_media FROM taxas_servico ORDER BY servico")
    rows = cursor.fetchall()
    conn.close()
    return rows
def update_servico(servico, taxa_media):
    conn = sqlite3.connect("taxas_servico.db")
    cursor = conn.cursor()
    cursor.execute("UPDATE taxas_servico SET taxa_media = ? WHERE servico = ?", (taxa_media, servico))
    conn.commit()
    conn.close()    
#========= DEFs PARA GERIR TAXAS ========== FIM ===== #


#========= DEFs PARA GERIR USERS ========== INICIO =====#
def tela_gerenciar_usuarios():
    st.title("Gerenciar Usuários")
    conn = sqlite3.connect("usuarios.db")
    cursor = conn.cursor()

    # Carrega DataFrame dos usuários
    df = pd.read_sql_query("SELECT id, nome, posicao, email, admin FROM usuarios", conn)
    st.dataframe(df)

    st.markdown("---")
    st.header("Cadastrar/Editar Usuário")

    opcoes = ["Novo Usuário"] + [
        f"{row['id']} - {row['nome']} ({row['email']})" for _, row in df.iterrows()
    ]
    escolha = st.selectbox("Selecione para editar:", opcoes, key="usuario_selecionado")

    # Valores default
    user_id = None
    nome = ""
    posicao = "Solicitante"
    email = ""
    admin = False

    # Posições possíveis (ajuste conforme suas opções originais)
    opcoes_posicao = ["Comercial", "Operação", "Elo Executivo"]

    # Se editar, preenche os campos
    if escolha != "Novo Usuário":
        selected_id = int(escolha.split(" - ")[0])
        st.session_state['user_id_edicao'] = selected_id
        user_row_df = df[df["id"] == selected_id]
        if not user_row_df.empty:
            user_row = user_row_df.iloc[0]
            user_id = user_row["id"]
            nome = user_row["nome"]
            posicao = user_row["posicao"]
            email = user_row["email"]
            admin = bool(user_row["admin"])
    else:
        st.session_state['user_id_edicao'] = None

    nome = st.text_input("Nome", value=nome)
    posicao = st.selectbox("Posição", opcoes_posicao, index=opcoes_posicao.index(posicao) if posicao in opcoes_posicao else 0)
    email = st.text_input("Email", value=email)
    senha = st.text_input("Senha (deixe em branco para não alterar)", type="password", key="senha_usuario")
    admin = st.checkbox("Administrador", value=admin)
    if user_id is not None:
        st.info(f"ID do usuário: {user_id}")

    col1, col2 = st.columns([2,1])
    with col1:
        if st.button("Salvar"):
            if senha:
                senha_hash = hashlib.sha256(senha.encode()).hexdigest()
            else:
                senha_hash = None

            if user_id is None:
                # Novo usuário
                if not senha_hash:
                    st.error("Informe uma senha para novo usuário!")
                else:
                    try:
                        cursor.execute(
                            "INSERT INTO usuarios (nome, posicao, email, senha, admin) VALUES (?, ?, ?, ?, ?)",
                            (nome, posicao, email, senha_hash, int(admin))
                        )
                        conn.commit()
                        st.success("Usuário cadastrado!")
                        #st.session_state["senha_usuario"] = ""  # Limpa o campo senha
                        st.rerun()
                    except sqlite3.IntegrityError:
                        st.error("E-mail já cadastrado.")
            else:
                # Edição de usuário
                if senha_hash:
                    cursor.execute(
                        "UPDATE usuarios SET nome=?, posicao=?, email=?, senha=?, admin=? WHERE id=?",
                        (nome, posicao, email, senha_hash, int(admin), user_id)
                    )
                else:
                    cursor.execute(
                        "UPDATE usuarios SET nome=?, posicao=?, email=?, admin=? WHERE id=?",
                        (nome, posicao, email, int(admin), user_id)
                    )
                conn.commit()
                st.success("Usuário atualizado!")
                #st.session_state["senha_usuario"] = ""  # Limpa o campo senha
                st.rerun()
    with col2:
        if user_id is not None and st.button("Resetar Senha"):
            # Reseta a senha do usuário para 'Inicial@123'
            senha_inicial = "Inicial@123"
            senha_hash = hashlib.sha256(senha_inicial.encode()).hexdigest()
            cursor.execute("UPDATE usuarios SET senha=? WHERE id=?", (senha_hash, user_id))
            conn.commit()
            st.success("Senha resetada para: Inicial@123")
            #st.session_state["senha_usuario"] = ""  # Limpa o campo senha
            st.rerun()

    # Exclusão protegida
    st.markdown("---")
    st.header("Excluir Usuário Selecionado")
    id_para_excluir = st.session_state.get('user_id_edicao', None)
    user_row_df = df[df["id"] == id_para_excluir] if id_para_excluir is not None else pd.DataFrame()
    if not user_row_df.empty:
        user_row = user_row_df.iloc[0]
        if user_row["admin"]:
            st.warning("Não é possível excluir um usuário Administrador! Remova o privilégio de admin antes.")
        else:
            if st.button("Excluir Usuário"):
                cursor.execute("DELETE FROM usuarios WHERE id=?", (id_para_excluir,))
                conn.commit()
                st.success("Usuário excluído!")
                st.rerun()
    else:
        st.info("Selecione um usuário para excluir.")

    conn.close()
#========= DEFs PARA GERIR USERS ========== FIM =====#

#========= DEFs PARA GERIR TEMPLATES ========== INICIO =====#
def tela_atualizar_modelos_pptx():
    st.title("Atualizar Modelos de Proposta (PPTX)")

    opcoes_modelo = {
        "AMS": "modelo_ams.pptx",
        "FÁBRICA": "modelo_fabrica.pptx",
        "ALOCAÇÃO": "modelo_aloc.pptx"
    }
    tipo = st.selectbox("Tipo de Proposta", list(opcoes_modelo.keys()))
    nome_modelo = opcoes_modelo[tipo]
    pasta_modelos = "Modelos"  # ajuste conforme necessário

    st.write(f"Modelo atual: `{nome_modelo}`")

    # Use um key dinâmico para forçar a limpeza do file_uploader
    if "upload_key" not in st.session_state:
        st.session_state["upload_key"] = 0

    uploaded_file = st.file_uploader(
        "Faça upload do novo modelo (.pptx)", 
        type=["pptx"],
        key=f"fileuploader_{st.session_state['upload_key']}"
    )

    # Só mostra botão se arquivo for carregado
    if uploaded_file is not None:
        if st.button("Atualizar modelo PPTX"):
            caminho_modelo = os.path.join(pasta_modelos, nome_modelo)

            # Backup do antigo, se existir
            if os.path.exists(caminho_modelo):
                data_hora = datetime.now().strftime("%Y%m%d_%H%M%S")
                caminho_bkp = os.path.join(
                    pasta_modelos, f"{nome_modelo.replace('.pptx', f'_bkp_{data_hora}.pptx')}"
                )
                shutil.move(caminho_modelo, caminho_bkp)
                st.info(f"Backup criado: `{os.path.basename(caminho_bkp)}`")

            # Salva novo modelo
            with open(caminho_modelo, "wb") as f:
                f.write(uploaded_file.getbuffer())
            st.success(f"Novo modelo `{nome_modelo}` salvo com sucesso!")

            # Limpa o upload para novo envio (muda o key)
            st.session_state["upload_key"] += 1
            st.rerun()

    # Lista arquivos na pasta (opcional)
    st.markdown("#### Arquivos na pasta de modelos:")
    arquivos = os.listdir(pasta_modelos)
    for arq in sorted(arquivos):
        st.write(f"- {arq}")
#========= DEFs PARA GERIR TEMPLATES ========== FIM =====#


# ======== FERRAMENTA PRINCIPAL ========
def propostator_tool(usuario=None):
    criar_banco()  # garante banco
    
    if usuario.admin:
        aba = st.sidebar.radio(
            "Navegar",
            ["Propostator AMS", "Administrar Rates", "Gerenciar Usuários", "Atualizar Modelos Proposta"]
        )
        st.sidebar.success("Acesso: Administrador") 
    else:
        aba = "Propostator AMS"
        st.sidebar.info("Acesso: Solicitante")
        
    if aba == "Propostator AMS":
        # Limpa sessão de campos se pediu limpar
        if st.session_state.get("limpar_campos", False):
            for k in list(st.session_state.keys()):
                if k.startswith("campo_"):
                    del st.session_state[k]
            st.session_state["limpar_campos"] = False

        if usuario:
            if st.button("Logout"):
                st.session_state["logado"] = False
                st.session_state["tela"] = "login"
                st.session_state["usuario_logado"] = None
                st.rerun()
                return
                
            #if usuario.admin:
            #    st.sidebar.success("Acesso: Administrador")                    
            #else:
            #    st.sidebar.info("Acesso: Solicitante")
                
        st.title("PROPOSTATOR AMS")

        # --- CAMPO TIPO ---
        tipo = st.selectbox("Tipo de Proposta", ["AMS", "FÁBRICA", "ALOCAÇÃO"], key="campo_tipo")

        comercial = st.text_input("Comercial (nome e sobrenome)", key="campo_comercial")
        email = st.text_input("E-mail Numen", value=usuario.email if usuario else "", disabled=True, key="campo_email")
        cliente = st.text_input("Nome do Cliente", key="campo_cliente")

        OPCAO_O = f"'O' {cliente.replace(' ', '_')}"
        OPCAO_A = f"'A' {cliente.replace(' ', '_')}"
        opcoes_genero = ["", OPCAO_O, OPCAO_A]
        Genero_Cliente = st.selectbox("Gênero Cliente", opcoes_genero, key="campo_genero_cliente")

        # OPP só números e máx 8
        opp_raw = st.text_input("Número da Oportunidade", max_chars=8, key="campo_opp")
        opp = "".join([c for c in opp_raw if c.isdigit()])[:8]
        if opp != opp_raw:
            st.session_state["campo_opp"] = opp

        # --- CAMPOS DINÂMICOS ---
        baseline = ""
        frentes = ""
        data = ""
        objetivo = ""
        horas_mensais = ""
        tera_baseline = ""
        nivel = ""

        if tipo == "AMS":
            frentes = st.text_input("Frentes de Atendimento", key="campo_frentes")
            baseline = st.text_input("Baseline em Horas", key="campo_baseline")
            data = st.date_input("Data", key="campo_data").strftime("%d/%m/%Y")
            objetivo = st.text_area("Objetivo", key="campo_objetivo")

        elif tipo == "FÁBRICA":
            frentes = st.text_input("Frentes de Atendimento", key="campo_frentes")
            tera_baseline = st.selectbox("Terá Baseline?", ["Sim", "Não"], key="campo_tera_baseline")
            if tera_baseline == "Sim":
                baseline = st.text_input("Baseline em Horas", key="campo_baseline")
            data = st.date_input("Data", key="campo_data").strftime("%d/%m/%Y")
            objetivo = st.text_area("Objetivo", key="campo_objetivo")

        elif tipo == "ALOCAÇÃO":
            frentes = st.text_input("Frente de Atendimento", key="campo_frentes")
            nivel = st.selectbox("Nível", ["K1", "K2", "K3", "K4", "K5"], key="campo_nivel")
            horas_mensais = st.text_input("Horas Mensais", key="campo_horas_mensais")
            data = st.date_input("Data", key="campo_data").strftime("%d/%m/%Y")
            objetivo = st.text_area("Objetivo", key="campo_objetivo")

        # --- MODELO POR TIPO ---
        nome_arquivo_modelo = {
            "AMS": "modelo_ams.pptx",
            "FÁBRICA": "modelo_fabrica.pptx",
            "ALOCAÇÃO": "modelo_aloc.pptx"
        }[tipo]
        pptx_modelo = os.path.join("Modelos", nome_arquivo_modelo)

        # --- CAMPOS OBRIGATÓRIOS ---
        campos_obrigatorios = {"Cliente": cliente, "Número da Oportunidade": opp}
        if tipo == "AMS":
            campos_obrigatorios["Baseline em Horas"] = baseline
        elif tipo == "FÁBRICA" and tera_baseline == "Sim":
            campos_obrigatorios["Baseline em Horas"] = baseline
        elif tipo == "ALOCAÇÃO":
            campos_obrigatorios["Horas Mensais"] = horas_mensais
            campos_obrigatorios["Nível"] = nivel
        campos_faltando = [campo for campo, valor in campos_obrigatorios.items() if not valor.strip()]

        # --- BUSCA TAXAS DINÂMICAS ---
        taxa = 220.00  # padrão, mas sobrescreve abaixo
        if tipo == "ALOCAÇÃO" and nivel:
            # Busca no rate_aloc
            try:
                conn = sqlite3.connect("rate_aloc.db")
                c = conn.cursor()
                c.execute("SELECT taxa FROM rate_aloc WHERE nivel = ?", (nivel,))
                res = c.fetchone()
                if res:
                    taxa = float(res[0])
                conn.close()
            except Exception as e:
                st.warning("Erro ao buscar taxa de alocação. Usando padrão 220.")

        elif tipo in ["AMS", "FÁBRICA"]:
            try:
                conn = sqlite3.connect("taxas_servico.db")
                c = conn.cursor()
                c.execute("SELECT taxa_media FROM taxas_servico WHERE servico = ?", (tipo,))
                res = c.fetchone()
                if res:
                    taxa = float(res[0])
                conn.close()
            except Exception as e:
                st.warning("Erro ao buscar taxa de serviço. Usando padrão 220.")

        # --- CÁLCULO VALOR ---
        valor_monetario = ""
        try:
            if tipo == "ALOCAÇÃO" and horas_mensais.strip():
                resultado = float(horas_mensais.replace(",", ".")) * taxa
            elif tipo == "AMS" and baseline.strip():
                resultado = float(baseline.replace(",", ".")) * taxa
            elif tipo == "FÁBRICA" and tera_baseline == "Sim" and baseline.strip():
                resultado = float(baseline.replace(",", ".")) * taxa
            else:
                resultado = None
            if resultado is not None:
                valor_monetario = "R$ {:,.2f}".format(resultado).replace(",", "X").replace(".", ",").replace("X", ".")
                st.write(f"Valor Mensal da Proposta: {valor_monetario}")
        except ValueError:
            st.warning("Digite um valor numérico válido para baseline ou horas.")

        # --- SUBSTITUIÇÃO E GERAR ---
        def substituir_campos(pptx_modelo, substituicoes, saida_nome):
            prs = Presentation(pptx_modelo)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                for marcador, valor in substituicoes.items():
                                    if marcador in run.text:
                                        run.text = run.text.replace(marcador, str(valor))
            prs.save(saida_nome)

        col1, col2 = st.columns([2, 1])
        with col1:
            gerar = st.button("Gerar apresentação")
        with col2:
            limpar = st.button("Limpar campos")

        if limpar:
            st.session_state["limpar_campos"] = True
            st.rerun()

        if gerar:
            if campos_faltando:
                st.error(f"Preencha os seguintes campos obrigatórios: {', '.join(campos_faltando)}")
            elif not os.path.exists(pptx_modelo):
                st.error(f"Arquivo modelo não encontrado: {pptx_modelo}")
            else:
                substituicoes = {
                    "{{CLIENTE}}": cliente,
                    "{{DATA}}": data,
                    "{{FRENTES}}": frentes,
                    "{{VALOR}}": valor_monetario,
                    "{{OPP}}": opp,
                    "{{GENERO_CLIENTE}}": Genero_Cliente,
                    "{{BASELINE}}": baseline,
                    "{{COMERCIAL}}": comercial,
                    "{{EMAIL}}": email,
                    "{{OBJETIVO}}": objetivo,
                    "{{HORAS_MENSAIS}}": horas_mensais,
                    "{{TERA_BASELINE}}": tera_baseline,
                    "{{NIVEL}}": nivel
                }
                saida_nome = f"{tipo}_PT_{opp}_{cliente.replace(' ', '_')}_{data.replace('/', '-')}.pptx"
                substituir_campos(pptx_modelo, substituicoes, saida_nome)
                with open(saida_nome, "rb") as file:
                    st.success("Apresentação gerada com sucesso!")
                    st.download_button(
                        label="Baixar apresentação personalizada",
                        data=file,
                        file_name=saida_nome,
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                    )
                # --- Salva LOG no banco ---
                salvar_proposta_log(
                    usuario, tipo, comercial, email, cliente, Genero_Cliente, opp,
                    frentes, baseline, data, objetivo, valor_monetario, horas_mensais, tera_baseline, nivel
                )

        # --- HISTÓRICO ---
        mostrar_historico(usuario)
        
    elif aba == "Administrar Rates":
        st.set_page_config(page_title="Gerenciar Taxas", layout="centered")
        st.title("Gestão de Taxas - Propostator")
        
        st.header("Taxas por Nível de Alocação")
        aloc_rows = get_alocacoes()
        cols = st.columns(5)
        edited = False

        for idx, (nivel, taxa) in enumerate(aloc_rows):
            with cols[idx % 5]:
                nova_taxa = st.number_input(f"Taxa {nivel}", min_value=0.0, max_value=10000.0, value=float(taxa), step=1.0, key=f"aloc_{nivel}")
                if nova_taxa != taxa:
                    if st.button(f"Salvar {nivel}"):
                        update_aloc(nivel, nova_taxa)
                        st.success(f"Taxa de {nivel} atualizada para R$ {nova_taxa:.2f}")
                        edited = True

        if edited:
            st.rerun()

        st.markdown("---")
        st.header("Taxas por Serviço")
        servico_rows = get_servicos()
        cols2 = st.columns(len(servico_rows))
        edited2 = False

        for idx, (servico, taxa_media) in enumerate(servico_rows):
            with cols2[idx]:
                nova_taxa = st.number_input(f"Taxa {servico}", min_value=0.0, max_value=10000.0, value=float(taxa_media), step=1.0, key=f"servico_{servico}")
                if nova_taxa != taxa_media:
                    if st.button(f"Salvar {servico}"):
                        update_servico(servico, nova_taxa)
                        st.success(f"Taxa de {servico} atualizada para R$ {nova_taxa:.2f}")
                        edited2 = True

        if edited2:
            st.rerun()

        st.info("Pronto! Todas as taxas podem ser atualizadas rapidamente por aqui.")
    
    elif aba == "Gerenciar Usuários":
        if not (usuario and getattr(usuario, "admin", 0)):
            st.warning("Acesso restrito ao administrador.")
            st.stop()
        tela_gerenciar_usuarios()

    elif aba == "Atualizar Modelos Proposta":
        tela_atualizar_modelos_pptx()
